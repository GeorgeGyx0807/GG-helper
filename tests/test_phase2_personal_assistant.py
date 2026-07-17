import sys
import zipfile
import json
from types import SimpleNamespace
from pathlib import Path

import pytest

from poppy.application.controller import DesktopController
from poppy.features.document_index import DocumentIndex
from poppy.features.document_extractors import extract_document
from poppy.features import document_extractors
from poppy.integrations import macos
from poppy.storage import AppPaths, DesktopDatabase
from poppy.tools import validate_tool
from poppy.tool_context import ToolContext


def test_library_index_is_grant_bounded_and_revocation_is_immediate(tmp_path):
    source_root = tmp_path / "notes"
    source_root.mkdir()
    note = source_root / "meeting.md"
    note.write_text("Project decision: ship Poppy in phases.\n", encoding="utf-8")
    paths = AppPaths(tmp_path / "app-data").ensure()
    database = DesktopDatabase(paths.database)
    grant = database.add_grant(source_root, can_read=True)
    index = DocumentIndex(database)
    source = index.add_source(source_root, grant)

    assert index.reindex(source["id"])[0]["documents"] == 1
    assert index.search("ship Poppy", [grant])[0]["path"] == str(note.resolve())

    database.delete_grant(grant["id"])
    assert index.search("ship Poppy", database.list_grants()) == []


def test_library_search_relaxes_multi_term_queries_and_reports_document_count(tmp_path):
    source_root = tmp_path / "notes"
    source_root.mkdir()
    note = source_root / "poppy-context.md"
    note.write_text("Poppy 的上下文由 ContextManager 管理，并在窗口超限时压缩历史。\n", encoding="utf-8")
    database = DesktopDatabase(AppPaths(tmp_path / "data").ensure().database)
    grant = database.add_grant(source_root, can_read=True)
    index = DocumentIndex(database)
    source = index.add_source(source_root, grant)

    index.reindex(source["id"])
    hits = index.search("poppy 上下文 context 设计", [grant])

    assert hits[0]["path"] == str(note.resolve())
    assert hits[0]["match_score"] >= 2
    assert database.get_library_source(source["id"])["document_count"] == 1


def test_library_search_expands_chinese_document_questions_to_english_terms(tmp_path):
    source_root = tmp_path / "notes"
    source_root.mkdir()
    note = source_root / "architecture.md"
    note.write_text(
        "Poppy uses a layered memory architecture. Session context is restored from checkpoints.\n",
        encoding="utf-8",
    )
    database = DesktopDatabase(AppPaths(tmp_path / "data").ensure().database)
    grant = database.add_grant(source_root, can_read=True)
    index = DocumentIndex(database)
    source = index.add_source(source_root, grant)

    index.reindex(source["id"])
    hits = index.search("这个文件的记忆和上下文是怎么设计的？", [grant])

    assert hits
    assert hits[0]["path"] == str(note.resolve())
    assert "memory architecture" in hits[0]["content"]


def test_locate_selection_normalizes_pdf_line_breaks_and_returns_neighbors(tmp_path):
    source_root = tmp_path / "papers"
    source_root.mkdir()
    paper = source_root / "2-poppy.txt"
    paper.write_text(
        "Introduction and motivation.\n"
        "The memory architecture main-\n"
        "tains session context in bounded checkpoints for reliable recovery.\n"
        "The next section discusses evaluation.\n",
        encoding="utf-8",
    )
    database = DesktopDatabase(AppPaths(tmp_path / "data").ensure().database)
    grant = database.add_grant(source_root, can_read=True)
    index = DocumentIndex(database)
    source = index.add_source(source_root, grant)
    index.reindex(source["id"])

    match = index.locate_selection(
        "The memory architecture maintains session context in bounded checkpoints",
        [grant],
        window_title="2-poppy.txt",
    )

    assert match is not None
    assert match["display_name"] == "2-poppy.txt"
    assert match["confidence"] >= 0.6
    assert "memory architecture" in match["context_rows"][0]["content"].lower()


def test_locate_selection_uses_filename_hint_and_rejects_low_confidence(tmp_path):
    source_root = tmp_path / "papers"
    source_root.mkdir()
    first = source_root / "first-paper.txt"
    second = source_root / "target-paper.txt"
    shared = "The layered memory architecture restores session context from bounded checkpoints.\n"
    first.write_text(shared, encoding="utf-8")
    second.write_text(shared, encoding="utf-8")
    database = DesktopDatabase(AppPaths(tmp_path / "data").ensure().database)
    grant = database.add_grant(source_root, can_read=True)
    index = DocumentIndex(database)
    source = index.add_source(source_root, grant)
    index.reindex(source["id"])

    match = index.locate_selection(
        "The layered memory architecture restores session context from bounded checkpoints.",
        [grant],
        window_title="target-paper.pdf",
    )

    assert match is not None
    assert match["display_name"] == "target-paper.txt"
    assert index.locate_selection("unrelated speculative sentence with no supporting evidence", [grant]) is None


def test_search_document_is_limited_to_one_authorized_document(tmp_path):
    source_root = tmp_path / "papers"
    source_root.mkdir()
    target = source_root / "target.md"
    distractor = source_root / "distractor.md"
    target.write_text("Target eviction uses a grace period before permanent removal.\n", encoding="utf-8")
    distractor.write_text("Distractor eviction uses no grace period and must never leak.\n", encoding="utf-8")
    database = DesktopDatabase(AppPaths(tmp_path / "data").ensure().database)
    grant = database.add_grant(source_root, can_read=True)
    index = DocumentIndex(database)
    source = index.add_source(source_root, grant)
    index.reindex(source["id"])
    target_document = database.get_document_by_path(target)

    hits = index.search_document("eviction grace period", target_document["id"], [grant])

    assert hits
    assert {row["path"] for row in hits} == {str(target.resolve())}
    database.delete_grant(grant["id"])
    assert index.search_document("eviction grace period", target_document["id"], []) == []


def test_authorizing_a_folder_makes_it_searchable_from_normal_chat(tmp_path):
    source_root = tmp_path / "authorized"
    source_root.mkdir()
    (source_root / "facts.md").write_text("Poppy 自动索引授权文件夹。\n", encoding="utf-8")
    paths = AppPaths(tmp_path / "app-data").ensure()
    controller = DesktopController(paths=paths, database=DesktopDatabase(paths.database))

    grant = controller.add_grant(source_root, can_read=True)

    assert controller.list_library_sources()[0]["document_count"] == 1
    assert controller.search_library("自动索引", limit=5)[0]["path"].endswith("facts.md")
    assert grant["path"] == str(source_root.resolve())


def test_library_source_cannot_escape_authorized_grant(tmp_path):
    allowed = tmp_path / "allowed"
    outside = tmp_path / "outside"
    allowed.mkdir()
    outside.mkdir()
    database = DesktopDatabase(AppPaths(tmp_path / "data").ensure().database)
    grant = database.add_grant(allowed, can_read=True)
    with pytest.raises(PermissionError):
        DocumentIndex(database).add_source(outside, grant)


def test_system_tool_validation_rejects_unsafe_url_and_accepts_library_query(tmp_path):
    context = ToolContext(
        root=Path(tmp_path),
        path_resolver=lambda value: Path(tmp_path) / value,
        shell_env_provider=lambda: {},
        depth=0,
        max_depth=1,
        spawn_delegate=lambda _args: "",
        library_searcher=lambda query, limit: [],
    )
    with pytest.raises(ValueError):
        validate_tool(context, "web_read", {"url": "javascript:alert(1)"})
    validate_tool(context, "library_search", {"query": "会议", "limit": 5})


def test_mac_integration_date_and_url_validation_are_side_effect_free():
    with pytest.raises(ValueError):
        macos._validate_url("file:///etc/passwd")
    assert macos._parse_calendar_date("2026-07-14T10:00:00").year == 2026


def test_system_approval_rules_are_scoped_to_tool_and_domain(tmp_path):
    paths = AppPaths(tmp_path / "data").ensure()
    controller = DesktopController(paths=paths, database=DesktopDatabase(paths.database))
    request = {"tool_name": "browser_open", "arguments": {"url": "https://example.com/docs"}}
    descriptor = controller._approval_descriptor(request)
    assert descriptor == ("browser_open", "system", "example.com")
    assert controller._approval_rule_allows(request) is False
    assert controller._save_approval_rule(request) is True
    assert controller._approval_rule_allows(request) is True


def test_docx_and_xlsx_extractors_return_searchable_text_and_locations(tmp_path):
    xlsx = tmp_path / "budget.xlsx"
    import openpyxl

    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet.title = "预算"
    sheet.append(["项目", "金额"])
    sheet.append(["Poppy 文档索引", 128])
    workbook.save(xlsx)
    workbook.close()

    spreadsheet = extract_document(xlsx, set())
    assert "Poppy 文档索引" in spreadsheet.content
    assert spreadsheet.chunks[0]["location"] == {
        "kind": "spreadsheet",
        "sheet": "预算",
        "row_start": 1,
        "row_end": 2,
    }

    docx = tmp_path / "notes.docx"
    with zipfile.ZipFile(docx, "w") as archive:
        archive.writestr(
            "[Content_Types].xml",
            """<?xml version='1.0' encoding='UTF-8'?>
            <Types xmlns='http://schemas.openxmlformats.org/package/2006/content-types'>
              <Default Extension='rels' ContentType='application/vnd.openxmlformats-package.relationships+xml'/>
              <Default Extension='xml' ContentType='application/xml'/>
              <Override PartName='/word/document.xml' ContentType='application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml'/>
            </Types>""",
        )
        archive.writestr(
            "_rels/.rels",
            """<Relationships xmlns='http://schemas.openxmlformats.org/package/2006/relationships'>
              <Relationship Id='rId1' Type='http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument' Target='word/document.xml'/>
            </Relationships>""",
        )
        archive.writestr(
            "word/document.xml",
            """<?xml version='1.0' encoding='UTF-8' standalone='yes'?>
            <w:document xmlns:w='http://schemas.openxmlformats.org/wordprocessingml/2006/main'>
              <w:body><w:p><w:r><w:t>Poppy Word 索引测试</w:t></w:r></w:p></w:body>
            </w:document>""",
        )

    word = extract_document(docx, set())
    assert "Poppy Word 索引测试" in word.content
    assert word.chunks[0]["location"]["kind"] == "docx"


def test_pptx_extractor_returns_searchable_slide_text(tmp_path):
    from pptx import Presentation

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Poppy 演示文稿"
    slide.placeholders[1].text = "统一文档问答支持 PowerPoint"
    pptx_path = tmp_path / "briefing.pptx"
    presentation.save(pptx_path)

    extracted = extract_document(pptx_path, set())

    assert "统一文档问答" in extracted.content
    assert extracted.chunks[0]["location"]["kind"] == "pptx"
    assert extracted.extractor == "markitdown-pptx"


def test_document_extraction_cache_is_invalidated_only_when_file_changes(tmp_path, monkeypatch):
    pdf = tmp_path / "cached.pdf"
    pdf.write_bytes(b"first")
    calls = []

    def fake_extract(path):
        calls.append(path.read_bytes())
        return document_extractors.ExtractedDocument("缓存内容", [{
            "line_start": 1,
            "line_end": 1,
            "content": "缓存内容",
            "location": {"kind": "pdf_page", "page": 1},
        }], "fake-pdf")

    document_extractors._extract_document_cached.cache_clear()
    monkeypatch.setattr(document_extractors, "_extract_pdf", fake_extract)

    assert extract_document(pdf, set()).content == "缓存内容"
    assert extract_document(pdf, set()).content == "缓存内容"
    assert calls == [b"first"]

    pdf.write_bytes(b"second version")
    assert extract_document(pdf, set()).content == "缓存内容"
    assert calls == [b"first", b"second version"]


def test_image_only_pdf_uses_native_ocr_fallback(tmp_path, monkeypatch):
    from PIL import Image

    pdf = tmp_path / "scan.pdf"
    Image.new("RGB", (100, 100), "white").save(pdf, "PDF")
    document_extractors._extract_document_cached.cache_clear()
    monkeypatch.setattr(document_extractors, "_native_ocr_helper", lambda: Path("/bin/true"))
    monkeypatch.setattr(
        document_extractors.subprocess,
        "run",
        lambda *_args, **_kwargs: SimpleNamespace(
            stdout=json.dumps(["扫描页中的中文记忆内容"], ensure_ascii=False)
        ),
    )

    extracted = extract_document(pdf, set())

    assert extracted.extractor == "macos-vision-ocr"
    assert "中文记忆" in extracted.content
    assert extracted.chunks[0]["location"]["kind"] == "pdf_page"
    assert extracted.chunks[0]["location"]["page"] == 1


def test_pdf_extractor_preserves_page_location_without_reading_outside_path(tmp_path, monkeypatch):
    class FakePage:
        def __init__(self, text):
            self.text = text

        def extract_text(self):
            return self.text

    class FakePdf:
        pages = [FakePage("第一页内容"), FakePage("第二页有会议关键词")]

        def __enter__(self):
            return self

        def __exit__(self, *_args):
            return False

    pdfplumber = SimpleNamespace(open=lambda _path: FakePdf())
    monkeypatch.setitem(sys.modules, "pdfplumber", pdfplumber)
    pdf = tmp_path / "meeting.pdf"
    pdf.write_bytes(b"test")

    extracted = extract_document(pdf, set())
    assert "会议关键词" in extracted.content
    assert extracted.chunks[1]["location"]["kind"] == "pdf_page"
    assert extracted.chunks[1]["location"]["page"] == 2


def test_library_search_returns_spreadsheet_locator(tmp_path):
    source_root = tmp_path / "library"
    source_root.mkdir()
    xlsx = source_root / "facts.xlsx"
    import openpyxl

    workbook = openpyxl.Workbook()
    workbook.active.append(["关键词", "个人资料库"])
    workbook.save(xlsx)
    workbook.close()

    database = DesktopDatabase(AppPaths(tmp_path / "data").ensure().database)
    grant = database.add_grant(source_root, can_read=True)
    index = DocumentIndex(database)
    source = index.add_source(source_root, grant)
    assert index.reindex(source["id"])[0]["documents"] == 1
    hit = index.search("个人资料库", [grant])[0]
    assert hit["location"]["kind"] == "spreadsheet"
    assert hit["location"]["sheet"] == "Sheet"


def test_reindex_reports_unextractable_pdf_instead_of_silently_indexing_it(tmp_path, monkeypatch):
    class EmptyPdf:
        pages = []

        def __enter__(self):
            return self

        def __exit__(self, *_args):
            return False

    monkeypatch.setitem(sys.modules, "pdfplumber", SimpleNamespace(open=lambda _path: EmptyPdf()))
    source_root = tmp_path / "library"
    source_root.mkdir()
    (source_root / "scan.pdf").write_bytes(b"not-text")
    database = DesktopDatabase(AppPaths(tmp_path / "data").ensure().database)
    grant = database.add_grant(source_root, can_read=True)
    index = DocumentIndex(database)
    source = index.add_source(source_root, grant)

    result = index.reindex(source["id"])[0]
    assert result["documents"] == 0
    assert result["errors"][0]["path"].endswith("scan.pdf")
